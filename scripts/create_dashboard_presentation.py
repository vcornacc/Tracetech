from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

OUT_PATH = Path(__file__).resolve().parents[1] / "Dashboard_TraceTech_Presentazione.pptx"

TITLE_COLOR = RGBColor(15, 23, 42)
SUBTITLE_COLOR = RGBColor(71, 85, 105)
ACCENT_COLOR = RGBColor(14, 116, 144)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_box = slide.shapes.title
    subtitle_box = slide.placeholders[1]

    title_box.text = title
    subtitle_box.text = subtitle

    title_run = title_box.text_frame.paragraphs[0].runs[0]
    title_run.font.size = Pt(40)
    title_run.font.bold = True
    title_run.font.color.rgb = TITLE_COLOR

    sub_run = subtitle_box.text_frame.paragraphs[0].runs[0]
    sub_run.font.size = Pt(20)
    sub_run.font.color.rgb = SUBTITLE_COLOR


def add_bullets_slide(prs: Presentation, title: str, bullets: list[str], note: str | None = None) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    t_run = slide.shapes.title.text_frame.paragraphs[0].runs[0]
    t_run.font.color.rgb = TITLE_COLOR
    t_run.font.bold = True

    body = slide.shapes.placeholders[1].text_frame
    body.clear()

    for i, bullet in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(30, 41, 59)

    if note:
        p = body.add_paragraph()
        p.text = ""
        p = body.add_paragraph()
        p.text = note
        p.level = 0
        p.font.size = Pt(18)
        p.font.italic = True
        p.font.color.rgb = ACCENT_COLOR


def build_presentation() -> Path:
    prs = Presentation()

    add_title_slide(
        prs,
        "TraceTech Dashboard",
        "Spiegazione semplice per chi parte da zero",
    )

    add_bullets_slide(
        prs,
        "Obiettivo del Dashboard",
        [
            "Far capire subito lo stato dei materiali critici.",
            "Mostrare i rischi principali in modo visuale e chiaro.",
            "Aiutare a decidere dove intervenire prima.",
            "Ridurre tempo perso tra file, report e sistemi diversi.",
        ],
        "In una frase: un cruscotto unico per monitorare rischio e recupero.",
    )

    add_bullets_slide(
        prs,
        "Come e organizzato",
        [
            "Analytics: Dashboard, CRM Materials, BOM & Risk.",
            "Operations: ECU Inventory, Decision Engine, Simulation.",
            "Strategy: Executive Dashboard, Financial Engine, HaaS Readiness.",
            "Ogni sezione risponde a una domanda precisa: cosa succede, perche, cosa fare.",
        ],
    )

    add_bullets_slide(
        prs,
        "Dashboard principale (home)",
        [
            "4 KPI chiave: materiali tracciati, esposizione alta, recovery rate, rischio portfolio.",
            "Criticality Matrix: confronto tra criticita e impatto economico.",
            "Cluster Distribution: come sono distribuiti i materiali per categoria.",
            "Active Alerts: segnali attivi da gestire subito.",
        ],
    )

    add_bullets_slide(
        prs,
        "BOM & Risk (analisi file)",
        [
            "Caricamento drag & drop di CSV/XLSX della distinta base.",
            "Matching automatico dei materiali e analisi rischio.",
            "Report con punteggio, driver principali e scenari.",
            "Serve per capire rischio reale prodotto per prodotto.",
        ],
    )

    add_bullets_slide(
        prs,
        "ECU Inventory e Alert",
        [
            "Vista componenti ECU monitorati nel sistema.",
            "Stato recupero, rischio e materiali coinvolti.",
            "Alert live con severita e numero ECU impattate.",
            "Utile per priorizzare interventi operativi.",
        ],
    )

    add_bullets_slide(
        prs,
        "Executive e Financial",
        [
            "Executive Dashboard: vista sintetica per management.",
            "Financial Engine: impatti economici e valore recuperabile.",
            "Supporta decisioni su budget, fornitori e mitigazioni.",
            "Trasforma dati tecnici in messaggi business.",
        ],
    )

    add_bullets_slide(
        prs,
        "Funzione nuova: Reset Test",
        [
            "Pulsante 'Reset Test' nella Dashboard principale.",
            "Azzera i dati a zero per test UI e verifiche di regressione.",
            "Pulsante 'Ripristina Dati' per tornare subito ai dati reali.",
            "Permette demo e test senza toccare dati sorgente.",
        ],
    )

    add_bullets_slide(
        prs,
        "Cosa abbiamo implementato (breve)",
        [
            "Motore MRF (Material Risk Factor) per rischio materiali.",
            "Parser BOM con analisi automatica e report rischio.",
            "Dashboard semplificata con sole sezioni utili/reali.",
            "Workflow CI/CD GitHub Pages corretto e stabilizzato.",
        ],
    )

    add_bullets_slide(
        prs,
        "Come leggere i dati in 30 secondi",
        [
            "1) Guarda i 4 KPI: capisci subito lo stato generale.",
            "2) Controlla Active Alerts: cosa richiede azione ora.",
            "3) Apri BOM & Risk: quali prodotti/materiali espongono piu rischio.",
            "4) Vai su Executive/Financial: traduci il rischio in impatto business.",
        ],
    )

    add_bullets_slide(
        prs,
        "Messaggio finale",
        [
            "Questo dashboard non e solo grafica: e uno strumento decisionale.",
            "Rende visibili i rischi prima che diventino costi.",
            "Permette test rapidi grazie alla modalita Reset Test.",
            "Obiettivo: decisioni piu veloci, piu semplici, piu sicure.",
        ],
    )

    prs.save(OUT_PATH)
    return OUT_PATH


if __name__ == "__main__":
    out = build_presentation()
    print(f"Presentation created: {out}")
